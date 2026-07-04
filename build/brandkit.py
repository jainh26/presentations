"""Panacea Biotec / NikoMom brand kit — reusable pptx helpers.
On-brand per Panacea_Biotec_Brand_Spec.md."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- palette ----
PINK    = RGBColor(0xD7, 0x1E, 0x83)
PURPLE  = RGBColor(0x57, 0x30, 0x7F)
SKY     = RGBColor(0x4E, 0xA6, 0xDC)
BLUE    = RGBColor(0x47, 0x75, 0xE7)
VIOLET  = RGBColor(0x89, 0x71, 0xE1)
ORANGE  = RGBColor(0xF7, 0x96, 0x31)
GREEN   = RGBColor(0x6B, 0x9F, 0x25)
CHAR    = RGBColor(0x45, 0x45, 0x51)
LGREY   = RGBColor(0xD8, 0xD9, 0xDC)
GREY    = RGBColor(0x8C, 0x8C, 0x8C)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
SOFT    = RGBColor(0xF4, 0xF2, 0xF7)   # soft lavender-grey card fill
SOFTPK  = RGBColor(0xFB, 0xEA, 0xF3)   # soft pink card fill
GREYBG  = RGBColor(0xF2, 0xF2, 0xF2)

TITLE_FONT = "Aptos Display"
BODY_FONT  = "Aptos"

EMU_W, EMU_H = 12192000, 6858000  # 13.33 x 7.5

def new_deck():
    p = Presentation()
    p.slide_width = Emu(EMU_W)
    p.slide_height = Emu(EMU_H)
    return p

def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def _set_font(run, size, color, bold=False, italic=False, font=BODY_FONT):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font

def rect(slide, x, y, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE, shadow=False):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    if shadow:
        el = sp._element.spPr
        ef = el.makeelement(qn('a:effectLst'), {}); el.append(ef)
        sh = ef.makeelement(qn('a:outerShdw'), {'blurRad':'55000','dist':'25000','dir':'5400000','rotWithShape':'0'})
        ef.append(sh)
        clr = sh.makeelement(qn('a:srgbClr'), {'val':'8A8A9A'}); sh.append(clr)
        a = clr.makeelement(qn('a:alpha'), {'val':'33000'}); clr.append(a)
    return sp

def txt(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        space_after=4, line_spacing=1.0, wrap=True):
    """runs: list of paragraphs; each paragraph is a list of (text,size,color,bold,italic,font) tuples."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    for m in ('left','right','top','bottom'):
        setattr(tf, f'margin_{m}', 0)
    first = True
    for para in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align; p.space_after = Pt(space_after); p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for seg in para:
            seg = list(seg)
            t = seg[0]; s = seg[1]; c = seg[2]
            b = seg[3] if len(seg) > 3 else False
            i = seg[4] if len(seg) > 4 else False
            f = seg[5] if len(seg) > 5 else BODY_FONT
            r = p.add_run(); r.text = t
            _set_font(r, s, c, b, i, f)
    return tb

def logo(slide, x=10.72, y=0.30):
    """Corporate endorsement lockup, top-right: 7N mark + Panacea Biotec + tagline."""
    # 7N pink rounded badge
    b = rect(slide, x, y, 0.62, 0.62, fill=PINK, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    b.adjustments[0] = 0.28
    tf = b.text_frame; tf.word_wrap = False
    for m in ('left','right','top','bottom'): setattr(tf, f'margin_{m}', 0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "7N"; _set_font(r, 20, WHITE, True, False, TITLE_FONT)
    # wordmark
    txt(slide, x+0.72, y+0.03, 1.85, 0.62,
        [[("Panacea Biotec", 13, PURPLE, True, False, TITLE_FONT)],
         [("BEST IN GOODNESS", 7.5, CHAR, True, False, BODY_FONT)]],
        anchor=MSO_ANCHOR.MIDDLE, space_after=1)

def footer(slide, n, date="July 2026"):
    txt(slide, 0.5, 7.06, 3.0, 0.3, [[(date, 9, GREY)]])
    txt(slide, 4.5, 7.06, 4.33, 0.3, [[("Confidential. Circulation Restricted.", 9, GREY)]], align=PP_ALIGN.CENTER)
    txt(slide, 12.0, 7.06, 0.8, 0.3, [[(str(n), 9, GREY)]], align=PP_ALIGN.RIGHT)

def title(slide, text, kicker=None):
    if kicker:
        txt(slide, 0.5, 0.34, 8.5, 0.3, [[(kicker.upper(), 11, PINK, True)]], space_after=0)
        txt(slide, 0.5, 0.62, 9.9, 0.9, [[(text, 26, PURPLE, True, False, TITLE_FONT)]])
    else:
        txt(slide, 0.5, 0.42, 9.9, 1.0, [[(text, 28, PINK, True, False, TITLE_FONT)]])

def divider(slide, x=0.5, y=1.5, w=12.33):
    rect(slide, x, y, w, 0.014, fill=CHAR)
    rect(slide, x, y+0.05, w, 0.014, fill=PURPLE)

def content_slide(prs, n, ttl, kicker=None, date="July 2026"):
    s = _blank(prs)
    logo(s); title(s, ttl, kicker); footer(s, n, date)
    return s

def section_slide(prs, num, ttl, sub, n, bg=PURPLE):
    s = _blank(prs)
    rect(s, 0, 0, 13.34, 7.5, fill=bg)
    # big ghost number
    txt(s, 0.5, 0.6, 6, 3.2, [[(num, 200, WHITE, True, False, TITLE_FONT)]], anchor=MSO_ANCHOR.TOP)
    n_el = s.shapes[-1]
    # accent bar
    rect(s, 0.6, 4.05, 1.4, 0.10, fill=PINK)
    txt(s, 0.6, 4.35, 11.5, 1.4, [[(ttl, 40, WHITE, True, False, TITLE_FONT)]])
    txt(s, 0.62, 5.55, 10.8, 1.2, [[(sub, 15, RGBColor(0xE7,0xDD,0xF0), False, True)]], line_spacing=1.15)
    txt(s, 12.0, 7.06, 0.8, 0.3, [[(str(n), 9, RGBColor(0xC9,0xB8,0xD8))]], align=PP_ALIGN.RIGHT)
    return s

def workshop_opener(prs, wnum, title, payoff, agenda, n, accent=PINK, date="July 2026"):
    """Chapter-style opener for each Workshop. White with colour rail + 5-step tracker."""
    s = _blank(prs)
    rect(s, 0, 0, 13.34, 7.5, fill=WHITE)
    rect(s, 0, 0, 0.28, 7.5, fill=accent)
    rect(s, 0.28, 0, 0.10, 7.5, fill=PURPLE)
    logo(s)
    # tracker
    steps = ["Understand","Define","Meaning","Translate","Finalise"]
    tx = 0.7; tw = 2.28
    for i, st in enumerate(steps):
        cur = (i+1 == wnum); done = (i+1 < wnum)
        fill = accent if cur else (SOFTPK if done else SOFT)
        tcol = WHITE if cur else (accent if done else GREY)
        c = rect(s, tx, 0.55, tw, 0.5, fill=fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE); c.adjustments[0]=0.16
        txt(s, tx, 0.55, tw, 0.5, [[(f"{i+1}  {st}", 10.5, tcol, cur or done)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        tx += tw + 0.09
    txt(s, 0.7, 1.55, 9.0, 0.4, [[(f"WORKSHOP {wnum} OF 5", 13, accent, True)]])
    txt(s, 0.7, 1.95, 11.8, 1.1, [[(title, 38, PURPLE, True, False, TITLE_FONT)]])
    txt(s, 0.72, 3.05, 11.5, 0.6, [[(payoff, 16, accent, False, True)]], line_spacing=1.1)
    # agenda card
    card(s, 0.7, 3.85, 11.9, 2.75, fill=SOFT)
    txt(s, 1.0, 4.05, 11.0, 0.4, [[("ON THE AGENDA",11,accent,True)]])
    half = (len(agenda)+1)//2
    colL = agenda[:half]; colR = agenda[half:]
    txt(s, 1.0, 4.5, 5.6, 2.0, [[("•  "+a,12,CHAR)] for a in colL], space_after=7, line_spacing=1.08)
    txt(s, 6.9, 4.5, 5.5, 2.0, [[("•  "+a,12,CHAR)] for a in colR], space_after=7, line_spacing=1.08)
    footer(s, n, date)
    return s

def chip(slide, x, y, w, text, color=PINK, tcolor=WHITE, h=0.34, size=10):
    c = rect(slide, x, y, w, h, fill=color, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    c.adjustments[0] = 0.5
    tf = c.text_frame; tf.word_wrap = False
    for m in ('left','right','top','bottom'): setattr(tf, f'margin_{m}', 0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text; _set_font(r, size, tcolor, True)
    return c

def card(slide, x, y, w, h, fill=SOFT, line=None, radius=0.06):
    c = rect(slide, x, y, w, h, fill=fill, line=line, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    try: c.adjustments[0] = radius
    except Exception: pass
    return c
